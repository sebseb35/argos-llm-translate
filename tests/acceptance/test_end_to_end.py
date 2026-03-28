from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from typer.testing import CliRunner

from local_translator.api import preview_file, translate_file, translate_text
from local_translator.cli import app
from tests.integration.document_fixture_builder import (
    create_docx_fixture,
    create_pdf_fixture,
    create_pptx_fixture,
    create_text_fixtures,
    create_xlsx_fixture,
)


runner = CliRunner()


def _fake_argos_translate(self, text: str) -> str:
    translated = text
    replacements = [
        ("# Titre", "# Title"),
        ("Bonjour équipe.", "Hello team."),
        ("Consultez ", "See "),
        (" et exécutez ", " and run "),
        (" pour la version ", " for version "),
        ("Le système utilise ", "The system uses "),
        ("La ", "The "),
        (" la ", " the "),
        (" du ", " for "),
        (" démarre", " starts"),
        (" après", " after"),
        (" validation", " validation"),
        ("Titre", "Title"),
        ("Bonjour", "Hello"),
        ("équipe", "team"),
        ("recette", "recipe"),
        ("lot", "batch"),
    ]
    for source, target in replacements:
        translated = translated.replace(source, target)
    return translated


def _passthrough_post_edit(self, source: str, translated: str, glossary: dict[str, str] | None = None) -> str:
    return translated


def _failing_post_edit(self, source: str, translated: str, glossary: dict[str, str] | None = None) -> str:
    raise RuntimeError("simulated llm failure")


def _rewriting_post_edit(self, source: str, translated: str, glossary: dict[str, str] | None = None) -> str:
    return translated.replace("acceptance testing", "recipe").replace("work package", "batch")


def _write_glossary(path: Path) -> Path:
    path.write_text(
        """
source_language: fr
target_language: en
entries:
  recette: acceptance testing
  lot: work package
""".strip(),
        encoding="utf-8",
    )
    return path


def _write_llm_model(path: Path) -> Path:
    path.write_text("stub gguf payload", encoding="utf-8")
    return path


def test_text_translation_argos_end_to_end(monkeypatch):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    result = translate_text("Bonjour équipe.", source_lang="fr", target_lang="en", engine="argos")

    assert result.translated_text == "Hello team."
    assert result.report.segment_count == 1
    assert "LT GLOSSARY TERM" not in result.translated_text
    assert "__LT_" not in result.translated_text


def test_text_translation_hybrid_end_to_end_allows_fallback(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)
    monkeypatch.setattr("local_translator.engines.llm_engine.LLMEngine.post_edit", _failing_post_edit)
    model_path = _write_llm_model(tmp_path / "model.gguf")

    result = translate_text(
        "Bonjour équipe.",
        source_lang="fr",
        target_lang="en",
        engine="hybrid",
        llm_model=model_path,
    )

    assert result.translated_text == "Hello team."
    assert result.report.segment_count == 1
    assert result.report.fallback_count == 1
    assert "LT GLOSSARY TERM" not in result.translated_text
    assert "__LT_" not in result.translated_text


@pytest.mark.parametrize("engine", ["argos", "hybrid"])
def test_glossary_enforcement_end_to_end(monkeypatch, tmp_path, engine: str):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)
    if engine == "hybrid":
        monkeypatch.setattr("local_translator.engines.llm_engine.LLMEngine.post_edit", _rewriting_post_edit)

    glossary_path = _write_glossary(tmp_path / "glossary.yaml")
    kwargs = {
        "content": "La recette du lot 2 démarre après validation.",
        "source_lang": "fr",
        "target_lang": "en",
        "engine": engine,
        "glossary": glossary_path,
    }
    if engine == "hybrid":
        kwargs["llm_model"] = _write_llm_model(tmp_path / "model.gguf")

    result = translate_text(**kwargs)

    assert "acceptance testing" in result.translated_text
    assert "work package" in result.translated_text
    assert "recipe" not in result.translated_text
    assert "batch" not in result.translated_text
    assert "LT GLOSSARY TERM" not in result.translated_text
    assert "__LT_" not in result.translated_text
    assert "  " not in result.translated_text


def test_markdown_file_translation_cli_end_to_end(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    fixtures_dir = tmp_path / "fixtures"
    create_text_fixtures(fixtures_dir)
    input_path = fixtures_dir / "sample.md"
    output_path = tmp_path / "output.md"
    glossary_path = _write_glossary(tmp_path / "glossary.yaml")

    result = runner.invoke(
        app,
        [
            "translate",
            str(input_path),
            "--from",
            "fr",
            "--to",
            "en",
            "--output",
            str(output_path),
            "--engine",
            "argos",
            "--glossary",
            str(glossary_path),
        ],
    )

    assert result.exit_code == 0
    assert output_path.exists()

    translated = output_path.read_text(encoding="utf-8")
    assert "# Release Notes" in translated
    assert "The **translation pipeline** supports offline processing." in translated
    assert "- Validate CPU-only execution" in translated
    assert "LT GLOSSARY TERM" not in translated
    assert "__LT_" not in translated
    assert "\n\n" in translated


def test_docx_file_translation_end_to_end(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    input_path = tmp_path / "sample.docx"
    output_path = tmp_path / "sample.translated.docx"
    create_docx_fixture(input_path)

    result = translate_file(
        input_path,
        output_path,
        source_lang="fr",
        target_lang="en",
        engine="argos",
    )

    assert result.output_path == output_path
    assert output_path.exists()

    document = Document(str(output_path))
    paragraph_texts = [paragraph.text for paragraph in document.paragraphs]
    assert "Project Brief" in paragraph_texts[0]
    assert "migration milestones for the billing service." in paragraph_texts[1]
    assert "Service Level Objective" in paragraph_texts[2]
    assert "incident response policy" in paragraph_texts[2]
    assert "LT GLOSSARY TERM" not in "\n".join(paragraph_texts)
    assert "__LT_" not in "\n".join(paragraph_texts)


def test_pptx_file_translation_end_to_end(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    input_path = tmp_path / "sample.pptx"
    output_path = tmp_path / "sample.translated.pptx"
    create_pptx_fixture(input_path)

    result = translate_file(
        input_path,
        output_path,
        source_lang="fr",
        target_lang="en",
        engine="argos",
    )

    assert result.output_path == output_path
    assert output_path.exists()

    presentation = Presentation(str(output_path))
    first_slide = presentation.slides[0]
    assert first_slide.shapes.title.text == "Roadmap Update"
    body_lines = [paragraph.text for paragraph in first_slide.placeholders[1].text_frame.paragraphs]
    assert "Q2 priorities:" in body_lines[0]
    assert "Improve latency" in body_lines[1]
    assert "Finalize compliance checklist" in body_lines[2]
    all_text = "\n".join(
        paragraph.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    )
    assert "LT GLOSSARY TERM" not in all_text
    assert "__LT_" not in all_text


def test_xlsx_file_translation_end_to_end(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    input_path = tmp_path / "sample.xlsx"
    output_path = tmp_path / "sample.translated.xlsx"
    create_xlsx_fixture(input_path)

    result = translate_file(
        input_path,
        output_path,
        source_lang="fr",
        target_lang="en",
        engine="argos",
    )

    assert result.output_path == output_path
    assert output_path.exists()

    workbook = load_workbook(str(output_path))
    overview = workbook["Overview"]
    risks = workbook["Risks"]
    assert overview["A1"].value == "Metric"
    assert overview["B3"].value == "Forecast remains conservative"
    assert overview["B2"].value == 125000
    assert overview["B8"].value is True
    assert risks["A2"].value == "Vendor delay"
    workbook_text = "\n".join(
        str(cell.value)
        for worksheet in workbook.worksheets
        for row in worksheet.iter_rows()
        for cell in row
        if cell.value is not None
    )
    assert "LT GLOSSARY TERM" not in workbook_text
    assert "__LT_" not in workbook_text


def test_pdf_file_translation_end_to_end(monkeypatch, tmp_path):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    input_path = tmp_path / "sample.pdf"
    output_path = tmp_path / "sample.translated.txt"
    create_pdf_fixture(input_path)

    result = translate_file(
        input_path,
        output_path,
        source_lang="fr",
        target_lang="en",
        engine="argos",
    )

    assert result.output_path == output_path
    assert output_path.exists()

    original_reader = PdfReader(str(input_path))
    assert "Text Native PDF Sample" in original_reader.pages[0].extract_text()

    translated = output_path.read_text(encoding="utf-8")
    assert "Text Native PDF Sample" in translated
    assert "Contains quarterly metrics and action items." in translated
    assert "API version v2.1 is in rollout." in translated
    assert translated.endswith("\n")
    assert "LT GLOSSARY TERM" not in translated
    assert "__LT_" not in translated


def test_protected_tokens_end_to_end(monkeypatch):
    monkeypatch.setattr("local_translator.engines.argos_engine.ArgosEngine.translate", _fake_argos_translate)

    content = "Consultez https://example.com et exécutez `python app.py --port 8080` pour la version v2.3.4."
    result = translate_text(content, source_lang="fr", target_lang="en", engine="argos")

    assert "https://example.com" in result.translated_text
    assert "`python app.py --port 8080`" in result.translated_text
    assert "v2.3.4" in result.translated_text
    assert "LT GLOSSARY TERM" not in result.translated_text
    assert "__LT_" not in result.translated_text


def test_preview_command_end_to_end(tmp_path):
    fixtures_dir = tmp_path / "fixtures"
    create_text_fixtures(fixtures_dir)
    input_path = fixtures_dir / "sample.md"

    result = runner.invoke(app, ["preview", str(input_path)])
    preview = preview_file(input_path)

    assert result.exit_code == 0
    assert "segments detected" in result.output
    assert preview.segment_count == 7
    assert preview.segment_count > 0
