from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from local_translator.config import RuntimeConfig
from local_translator.models.types import EngineMode
from local_translator.pipeline.translator import TranslationPipeline
from local_translator.reconstructors.registry import get_extractor
from tests.integration.document_fixture_builder import (
    create_blank_pdf_fixture,
    create_docx_fixture,
    create_pdf_fixture,
    create_pptx_fixture,
    create_xlsx_fixture,
)

FIXTURES_DIR = Path(__file__).resolve().parent.parent / "fixtures"
TRANSLATED_TAG = " [TRANSLATED]"


class FakeTranslator:
    def translate(self, text: str) -> str:
        if not text.strip():
            return text
        return f"{text}{TRANSLATED_TAG}"


def prepare_fixture_path(tmp_path: Path, file_name: str) -> Path:
    if file_name in {"sample.txt", "sample.md"}:
        return FIXTURES_DIR / file_name

    generated_dir = tmp_path / "generated-fixtures"
    generated_dir.mkdir(exist_ok=True)
    fixture_path = generated_dir / file_name

    if file_name == "sample.docx":
        create_docx_fixture(fixture_path)
    elif file_name == "sample.pptx":
        create_pptx_fixture(fixture_path)
    elif file_name == "sample.xlsx":
        create_xlsx_fixture(fixture_path)
    elif file_name == "sample.pdf":
        create_pdf_fixture(fixture_path)
    elif file_name == "blank.pdf":
        create_blank_pdf_fixture(fixture_path)
    else:  # pragma: no cover - guard for future updates
        raise ValueError(f"Unsupported fixture request: {file_name}")
    return fixture_path


def run_document_pipeline(input_path: Path, output_path: Path, monkeypatch) -> Path:
    cfg = RuntimeConfig(source_lang="fr", target_lang="en", engine_mode=EngineMode.ARGOS)
    pipeline = TranslationPipeline(cfg)
    monkeypatch.setattr(pipeline, "argos", FakeTranslator())

    extractor = get_extractor(input_path)
    extracted = extractor.extract(input_path)
    translated_segments = [pipeline.translate_text(segment).text for segment in extracted.segments]
    extractor.reconstruct(extracted, translated_segments, output_path)
    return output_path


def test_txt_pipeline_with_fixture(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.txt")
    output_path = tmp_path / "sample.translated.txt"

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()
    content = output_path.read_text(encoding="utf-8")
    assert "Quarterly Summary" in content
    assert TRANSLATED_TAG in content


def test_md_pipeline_preserves_markdown_structure(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.md")
    output_path = tmp_path / "sample.translated.md"

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()
    lines = output_path.read_text(encoding="utf-8").splitlines()
    assert any(line.startswith("# ") for line in lines)
    assert any(line.startswith("- ") for line in lines)
    assert any(TRANSLATED_TAG in line for line in lines if line.strip())


def test_docx_pipeline_with_fixture(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.docx")
    output_path = tmp_path / "sample.translated.docx"

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()
    translated_doc = Document(str(output_path))

    paragraphs = translated_doc.paragraphs
    paragraph_texts = [paragraph.text for paragraph in paragraphs]

    expected_heading = f"Project Brief{TRANSLATED_TAG}"
    expected_intro = (
        f"This document describes migration milestones for the billing service.{TRANSLATED_TAG}"
    )

    assert paragraphs[0].style.name.startswith("Heading")
    assert paragraph_texts[0] == expected_heading
    assert paragraph_texts[1] == expected_intro

    assert "Key term:" in paragraph_texts[2]
    assert "Service Level Objective" in paragraph_texts[2]
    assert "incident response policy" in paragraph_texts[2]

    key_phrases_in_order = [expected_heading, expected_intro, paragraph_texts[2]]
    ordered_indexes = [paragraph_texts.index(phrase) for phrase in key_phrases_in_order]
    assert ordered_indexes == sorted(ordered_indexes)

    formatted_paragraph = paragraphs[2]
    assert formatted_paragraph.runs[1].bold is True
    assert formatted_paragraph.runs[3].italic is True
    assert all(TRANSLATED_TAG in run.text for run in formatted_paragraph.runs if run.text.strip())

    assert len(translated_doc.tables) == 1
    table = translated_doc.tables[0]
    table_text = "\n".join(cell.text for row in table.rows for cell in row.cells)
    assert "Owner" in table_text
    assert "Status" in table_text
    assert "Platform Team" in table_text
    assert "On Track" in table_text

    expected_cells = [
        [f"Owner{TRANSLATED_TAG}", f"Status{TRANSLATED_TAG}"],
        [f"Platform Team{TRANSLATED_TAG}", f"On Track{TRANSLATED_TAG}"],
    ]
    actual_cells = [[cell.text for cell in row.cells] for row in table.rows]
    assert actual_cells == expected_cells


def test_pptx_pipeline_with_fixture(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.pptx")
    output_path = tmp_path / "sample.translated.pptx"

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()
    translated_presentation = Presentation(str(output_path))
    assert len(translated_presentation.slides) == 2

    first_slide = translated_presentation.slides[0]
    second_slide = translated_presentation.slides[1]

    assert first_slide.shapes.title.text == f"Roadmap Update{TRANSLATED_TAG}"

    body_lines = [paragraph.text for paragraph in first_slide.placeholders[1].text_frame.paragraphs]
    assert body_lines == [
        f"Q2 priorities:{TRANSLATED_TAG}",
        f"Improve latency{TRANSLATED_TAG}",
        f"Finalize compliance checklist{TRANSLATED_TAG}",
    ]

    owner_notes_shape = next(
        shape
        for shape in first_slide.shapes
        if shape.has_text_frame and shape.text_frame.paragraphs[0].text == f"Owner notes{TRANSLATED_TAG}"
    )
    assert owner_notes_shape.text_frame.paragraphs[1].text == f"Keep rollout deterministic{TRANSLATED_TAG}"

    table_shape = next(shape for shape in first_slide.shapes if shape.has_table)
    table = table_shape.table
    assert table.cell(0, 0).text == f"Region{TRANSLATED_TAG}"
    assert table.cell(0, 1).text == f"Status{TRANSLATED_TAG}"
    assert table.cell(1, 0).text == f"US-East{TRANSLATED_TAG}"
    assert table.cell(1, 1).text == ""

    assert second_slide.shapes.title.text == f"No body text slide{TRANSLATED_TAG}"


def test_xlsx_pipeline_preserves_structure_and_non_text_cells(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.xlsx")
    output_path = tmp_path / "sample.translated.xlsx"

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()

    original_workbook = load_workbook(str(input_path))
    translated_workbook = load_workbook(str(output_path))

    assert translated_workbook.sheetnames == original_workbook.sheetnames
    assert len(translated_workbook.worksheets) == len(original_workbook.worksheets)

    overview_original = original_workbook["Overview"]
    overview = translated_workbook["Overview"]

    assert overview.dimensions == overview_original.dimensions

    # Text cells are translated.
    assert overview["A1"].value == f"Metric{TRANSLATED_TAG}"
    assert overview["B3"].value == f"Forecast remains conservative{TRANSLATED_TAG}"
    assert overview["B6"].value == f"00123{TRANSLATED_TAG}"
    assert overview["B9"].value == f"  keep spacing{TRANSLATED_TAG}  "

    # Non-text cells remain unchanged.
    assert overview["B2"].value == 125000
    assert overview["B4"].value == "=SUM(B2:B3)"
    assert overview["B7"].value == overview_original["B7"].value
    assert overview["B7"].is_date is True
    assert overview["B8"].value is True
    assert overview["B10"].value == overview_original["B10"].value
    assert overview["B11"].value is None

    risks_original = original_workbook["Risks"]
    risks = translated_workbook["Risks"]

    assert risks.dimensions == risks_original.dimensions
    assert risks["A2"].value == f"Vendor delay{TRANSLATED_TAG}"
    assert risks["D5"].value == f"Mitigation notes{TRANSLATED_TAG}"


def test_pdf_text_pipeline_writes_translated_text_sidecar_output(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.pdf")
    output_path = tmp_path / "sample.translated.txt"

    reader = PdfReader(str(input_path))
    extracted = reader.pages[0].extract_text()
    assert "Text Native PDF Sample" in extracted

    run_document_pipeline(input_path, output_path, monkeypatch)

    assert output_path.exists()
    content = output_path.read_text(encoding="utf-8")
    assert content.endswith("\n")
    assert "Text Native PDF Sample" in content
    assert "Contains quarterly metrics" in content
    assert TRANSLATED_TAG in content


def test_pdf_pipeline_rejects_non_text_output_extension(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "sample.pdf")
    output_path = tmp_path / "sample.translated.pdf"

    with pytest.raises(ValueError, match="text sidecar outputs only"):
        run_document_pipeline(input_path, output_path, monkeypatch)


def test_pdf_pipeline_rejects_image_like_pdf(tmp_path, monkeypatch):
    input_path = prepare_fixture_path(tmp_path, "blank.pdf")
    output_path = tmp_path / "blank.translated.txt"

    with pytest.raises(ValueError, match="Only text-native PDFs are supported"):
        run_document_pipeline(input_path, output_path, monkeypatch)

    assert not output_path.exists()


def test_pdf_pipeline_rejects_invalid_pdf(tmp_path, monkeypatch):
    input_path = tmp_path / "broken.pdf"
    input_path.write_text("not a real pdf", encoding="utf-8")

    with pytest.raises(ValueError, match="Invalid PDF file"):
        run_document_pipeline(input_path, tmp_path / "broken.translated.txt", monkeypatch)
