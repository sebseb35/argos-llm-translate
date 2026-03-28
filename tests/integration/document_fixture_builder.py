from __future__ import annotations

from datetime import date
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject
from pptx import Presentation
from pptx.util import Inches


FIXTURE_TEXT = {
    "txt": "Quarterly Summary\nThe platform uptime reached 99.9%.\nAction items: verify API v2 rollout and budget forecast.",
    "md": (
        "# Release Notes\n\n"
        "The **translation pipeline** supports offline processing.\n\n"
        "- Validate CPU-only execution\n"
        "- Keep version v1.2.3 stable\n"
        "- Review KPI dashboard URL: https://example.com/kpi\n"
    ),
}


def create_text_fixtures(fixtures_dir: Path) -> None:
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    (fixtures_dir / "sample.txt").write_text(FIXTURE_TEXT["txt"], encoding="utf-8")
    (fixtures_dir / "sample.md").write_text(FIXTURE_TEXT["md"], encoding="utf-8")


def create_docx_fixture(path: Path) -> None:
    doc = Document()
    doc.add_heading("Project Brief", level=1)
    doc.add_paragraph("This document describes migration milestones for the billing service.")

    formatted = doc.add_paragraph()
    formatted.add_run("Key term: ")
    bold_run = formatted.add_run("Service Level Objective")
    bold_run.bold = True
    formatted.add_run(" and ")
    italic_run = formatted.add_run("incident response policy")
    italic_run.italic = True
    formatted.add_run(".")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Platform Team"
    table.cell(1, 1).text = "On Track"
    doc.save(path)


def create_pptx_fixture(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap Update"

    body = slide.placeholders[1].text_frame
    body.paragraphs[0].text = "Q2 priorities:"
    body.add_paragraph().text = "Improve latency"
    body.add_paragraph().text = "Finalize compliance checklist"

    notes_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(5.5), Inches(1.2))
    notes_box.text_frame.paragraphs[0].text = "Owner notes"
    notes_box.text_frame.add_paragraph().text = "Keep rollout deterministic"

    table_shape = slide.shapes.add_table(rows=2, cols=2, left=Inches(6.0), top=Inches(1.6), width=Inches(3.0), height=Inches(1.6))
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "US-East"
    table.cell(1, 1).text = ""

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    second_slide.shapes.title.text = "No body text slide"
    presentation.save(path)


def create_xlsx_fixture(path: Path) -> None:
    workbook = Workbook()
    overview = workbook.active
    overview.title = "Overview"

    overview["A1"] = "Metric"
    overview["B1"] = "Value"
    overview["A2"] = "Revenue"
    overview["B2"] = 125000
    overview["A3"] = "Comment"
    overview["B3"] = "Forecast remains conservative"
    overview["A4"] = "Total"
    overview["B4"] = "=SUM(B2:B3)"

    overview["A6"] = "ID"
    overview["B6"] = "00123"
    overview["A7"] = "Launch Date"
    overview["B7"] = date(2025, 1, 31)
    overview["A8"] = "Approved"
    overview["B8"] = True
    overview["A9"] = "Whitespace"
    overview["B9"] = "  keep spacing  "
    overview["A10"] = "Empty String"
    overview["B10"] = ""
    overview["A11"] = "Missing"
    overview["B11"] = None

    risks = workbook.create_sheet("Risks")
    risks["A1"] = "Risk"
    risks["B1"] = "Owner"
    risks["A2"] = "Vendor delay"
    risks["B2"] = "Ops"
    risks["A3"] = "Escalation needed"
    risks["B3"] = "Program Manager"

    # Sparse text cell to validate coordinate preservation in sparse sheets.
    risks["D5"] = "Mitigation notes"

    workbook.save(path)


def create_pdf_fixture(path: Path) -> None:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)

    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )

    stream = DecodedStreamObject()
    stream.set_data(
        b"BT\n/F1 18 Tf\n72 740 Td\n(Text Native PDF Sample) Tj\n"
        b"0 -24 Td\n(Contains quarterly metrics and action items.) Tj\n"
        b"0 -24 Td\n(API version v2.1 is in rollout.) Tj\nET\n"
    )
    page[NameObject("/Contents")] = writer._add_object(stream)

    with path.open("wb") as output:
        writer.write(output)


def create_blank_pdf_fixture(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as output:
        writer.write(output)
