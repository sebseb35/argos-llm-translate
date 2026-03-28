from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.text.text import TextFrame

from local_translator.extractors.base import BaseExtractor, ExtractedDocument



@dataclass(slots=True)
class _ShapeRef:
    path: list[int]
    shape: BaseShape


class PptxExtractor(BaseExtractor):
    suffixes = (".pptx",)

    _OBJECTS_KEY = "objects"

    @staticmethod
    def _iter_shape_refs(shapes: SlideShapes, path_prefix: list[int] | None = None) -> list[_ShapeRef]:
        refs: list[_ShapeRef] = []
        prefix = path_prefix or []

        for shape_index, shape in enumerate(shapes):
            shape_path = [*prefix, shape_index]
            refs.append(_ShapeRef(path=shape_path, shape=shape))

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                refs.extend(PptxExtractor._iter_shape_refs(shape.shapes, path_prefix=shape_path))

        return refs

    @staticmethod
    def _resolve_shape_by_path(shapes: SlideShapes, shape_path: list[int]) -> BaseShape | None:
        current_shapes: Any = shapes
        current_shape: BaseShape | None = None

        for depth, index in enumerate(shape_path):
            if index < 0 or index >= len(current_shapes):
                return None
            current_shape = current_shapes[index]

            if depth < len(shape_path) - 1:
                if not (current_shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(current_shape, "shapes")):
                    return None
                current_shapes = current_shape.shapes

        return current_shape

    @staticmethod
    def _translated_value(extracted: ExtractedDocument, translated_segments: list[str], index: int) -> str:
        if index < len(translated_segments):
            return translated_segments[index]
        if index < len(extracted.segments):
            return extracted.segments[index]
        return ""

    @staticmethod
    def _write_text_frame(text_frame: TextFrame, values: list[str]) -> None:
        if len(values) == len(text_frame.paragraphs):
            for paragraph, value in zip(text_frame.paragraphs, values):
                paragraph.text = value
            return

        text_frame.clear()
        if not values:
            return

        text_frame.paragraphs[0].text = values[0]
        for value in values[1:]:
            text_frame.add_paragraph().text = value

    def extract(self, file_path: Path) -> ExtractedDocument:
        prs = Presentation(str(file_path))
        segments: list[str] = []
        objects: list[dict[str, Any]] = []

        for slide_index, slide in enumerate(prs.slides):
            for shape_ref in self._iter_shape_refs(slide.shapes):
                shape = shape_ref.shape

                if shape.has_text_frame:
                    paragraph_indexes: list[int] = []
                    for paragraph in shape.text_frame.paragraphs:
                        segments.append(paragraph.text)
                        paragraph_indexes.append(len(segments) - 1)

                    objects.append(
                        {
                            "slide_index": slide_index,
                            "shape_path": shape_ref.path,
                            "type": "text_frame",
                            "paragraph_indexes": paragraph_indexes,
                        }
                    )

                if shape.has_table:
                    table_cell_indexes: list[list[int]] = []
                    for row in shape.table.rows:
                        row_indexes: list[int] = []
                        for cell in row.cells:
                            segments.append(cell.text)
                            row_indexes.append(len(segments) - 1)
                        table_cell_indexes.append(row_indexes)

                    objects.append(
                        {
                            "slide_index": slide_index,
                            "shape_path": shape_ref.path,
                            "type": "table",
                            "table_cell_indexes": table_cell_indexes,
                        }
                    )

        return ExtractedDocument(file_path=file_path, segments=segments, metadata={self._OBJECTS_KEY: objects})

    def reconstruct(self, extracted: ExtractedDocument, translated_segments: list[str], output_path: Path) -> None:
        prs = Presentation(str(extracted.file_path))
        objects_meta = extracted.metadata.get(self._OBJECTS_KEY, []) if extracted.metadata else []

        for object_meta in objects_meta:
            slide_index = object_meta.get("slide_index")
            shape_path = object_meta.get("shape_path")
            object_type = object_meta.get("type")

            if not isinstance(slide_index, int) or not isinstance(shape_path, list):
                continue
            if slide_index < 0 or slide_index >= len(prs.slides):
                continue

            shape = self._resolve_shape_by_path(prs.slides[slide_index].shapes, shape_path)
            if shape is None:
                continue

            if object_type == "text_frame" and shape.has_text_frame:
                paragraph_indexes = object_meta.get("paragraph_indexes", [])
                values = [
                    self._translated_value(extracted=extracted, translated_segments=translated_segments, index=index)
                    for index in paragraph_indexes
                ]
                self._write_text_frame(shape.text_frame, values)

            if object_type == "table" and shape.has_table:
                table_cell_indexes = object_meta.get("table_cell_indexes", [])
                for row, row_indexes in zip(shape.table.rows, table_cell_indexes):
                    for cell, index in zip(row.cells, row_indexes):
                        cell.text = self._translated_value(
                            extracted=extracted,
                            translated_segments=translated_segments,
                            index=index,
                        )

        prs.save(str(output_path))
